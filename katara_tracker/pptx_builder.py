"""Build a native PowerPoint (.pptx) deck from Client records.

Re-creates the Katara member-profile design (maroon header, portrait, the
two-column field list, the green status note) with python-pptx, and adds
section divider slides (Pending / Approved / Paid) plus an Analysis section
(KPIs + membership-type table). Unlike the .odp generator this draws slides
from scratch, so it is a faithful re-creation rather than a pixel clone.
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from .analytics import kpis, membership_breakdown, status_counts
from .odp_parser import Client
from .theme import GREEN, MAROON, RED, STATUSES, STATUS_COLORS

MAROON_C = RGBColor.from_string(MAROON)
GREEN_C = RGBColor.from_string(GREEN)
RED_C = RGBColor.from_string(RED)
WHITE_C = RGBColor.from_string("FFFFFF")
BLACK_C = RGBColor.from_string("111111")

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

ROWS = [
    ("Application Date :", "application_date"),
    ("Name :", "name"),
    ("Age :", "age"),
    ("Occupation :", "occupation"),
    ("Company :", "company"),
    ("Membership Plan :", "membership_plan"),
    ("Rate :", "rate"),
    ("Tag :", "tag"),
    ("Special Request :", "special_request"),
    ("CEC :", "cec"),
]


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _box(slide, x, y, w, h):
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))


def _text(slide, x, y, w, h, text, size=14, bold=False, color=BLACK_C,
          align=PP_ALIGN.LEFT, wrap=True, anchor=MSO_ANCHOR.TOP):
    tb = _box(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def _rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE

    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _member_slide(prs, client: Client, media_dir: str) -> None:
    s = _blank(prs)

    # Header band.
    _rect(s, 0, 0, 13.333, 1.0, MAROON_C)
    _text(s, 0.3, 0.25, 2.4, 0.6, "KATARA", 26, True, WHITE_C)
    _text(s, 2.75, 0.22, 4.0, 0.6, "MEMBER PROFILE", 28, True, WHITE_C)
    if client.app_id:
        _text(s, 9.4, 0.30, 3.6, 0.5, client.app_id, 16, True, WHITE_C,
              align=PP_ALIGN.RIGHT)

    # Status note.
    if client.status in ("Approved", "Paid"):
        badge = _rect(s, 6.95, 0.28, 1.7, 0.5, GREEN_C)
        tf = badge.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = client.status
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = WHITE_C

    if client.application_type:
        _text(s, 0.35, 1.1, 6.0, 0.4, "Application Type: %s" %
              client.application_type, 16, True, MAROON_C)

    # Portrait.
    placed = False
    if client.photo:
        cand = client.photo
        if not os.path.exists(cand) and media_dir:
            cand = os.path.join(media_dir, os.path.basename(client.photo))
        if os.path.exists(cand):
            try:
                s.shapes.add_picture(cand, Inches(0.6), Inches(1.75),
                                     Inches(2.7), Inches(2.65))
                placed = True
            except Exception:
                placed = False
    if not placed:
        ph = _rect(s, 0.6, 1.75, 2.7, 2.65, RGBColor.from_string("EEE8EC"))
        _text(s, 0.6, 2.9, 2.7, 0.4, "No photo", 14, False, MAROON_C,
              align=PP_ALIGN.CENTER)
        del ph

    _text(s, 0.6, 4.5, 3.2, 0.5, client.name, 18, True, MAROON_C)
    if client.mobile:
        _text(s, 0.6, 4.95, 3.2, 0.4, "Mobile Number : %s" % client.mobile,
              13, False, BLACK_C)

    # Two-column field list.
    y = 1.75
    for label, attr in ROWS:
        val = getattr(client, attr, "") or "—"
        _text(s, 3.7, y, 1.9, 0.4, label, 13, True, MAROON_C)
        if attr == "special_request":
            _text(s, 5.6, y, 7.2, 1.1, str(val), 11, False,
                  RED_C if val and val != "—" else BLACK_C)
            y += 0.75
        else:
            _text(s, 5.6, y, 7.2, 0.4, str(val), 13, False, BLACK_C)
            y += 0.49

    # Footer accent.
    _rect(s, 0, 7.36, 13.333, 0.14,
          RGBColor.from_string(STATUS_COLORS.get(client.status, MAROON)))


def _divider_slide(prs, title: str, subtitle: str) -> None:
    s = _blank(prs)
    _rect(s, 0, 0, 13.333, 7.5, MAROON_C)
    _text(s, 1.0, 2.7, 11.333, 1.2, title, 54, True, WHITE_C,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _rect(s, 5.4, 3.95, 2.5, 0.06, GREEN_C)
    _text(s, 1.0, 4.1, 11.333, 0.8, subtitle, 24, False,
          RGBColor.from_string("E8D8E2"), align=PP_ALIGN.CENTER)


def _analysis_slides(prs, clients: list[Client]) -> None:
    # KPI slide.
    s = _blank(prs)
    _rect(s, 0, 0, 13.333, 1.0, MAROON_C)
    _text(s, 0.4, 0.22, 12.5, 0.6, "ANALYSIS — Membership Summary", 28, True,
          WHITE_C)
    y = 1.4
    for label, value, is_money in kpis(clients):
        disp = "{:,} QAR".format(value) if is_money else str(value)
        col = GREEN_C if "PAID" in label else BLACK_C
        _text(s, 0.8, y, 7.5, 0.45, label, 16, True, MAROON_C)
        _text(s, 8.3, y, 4.2, 0.45, disp, 16, "PAID" in label, col)
        y += 0.55

    # Membership-type table slide.
    s = _blank(prs)
    _rect(s, 0, 0, 13.333, 1.0, MAROON_C)
    _text(s, 0.4, 0.22, 12.5, 0.6, "ANALYSIS — By Membership Type", 28, True,
          WHITE_C)
    rows = membership_breakdown(clients)
    headers = ["Membership Type", "Joiners", "Paid", "Amount Paid (QAR)",
               "Pipeline (QAR)"]
    nrows = len(rows) + 2  # header + data + totals
    tbl_shape = s.shapes.add_table(
        nrows, 5, Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.7)
    )
    table = tbl_shape.table
    table.columns[0].width = Inches(5.3)
    for i, w in enumerate((1.4, 1.2, 2.9, 1.7), start=1):
        table.columns[i].width = Inches(w)

    def setcell(r, c, text, bold=False, color=BLACK_C, align=PP_ALIGN.LEFT):
        cell = table.cell(r, c)
        cell.margin_left = cell.margin_right = Pt(4)
        cell.margin_top = cell.margin_bottom = Pt(1)
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(12)
        run.font.bold = bold
        run.font.color.rgb = color

    for c, h in enumerate(headers):
        setcell(0, c, h, True, WHITE_C, PP_ALIGN.CENTER)
        table.cell(0, c).fill.solid()
        table.cell(0, c).fill.fore_color.rgb = MAROON_C
    tot = {"n": 0, "paid": 0, "paid_amt": 0, "pipe": 0}
    for r, row in enumerate(rows, start=1):
        setcell(r, 0, row["type"])
        setcell(r, 1, str(row["n"]), align=PP_ALIGN.CENTER)
        setcell(r, 2, str(row["paid"]), align=PP_ALIGN.CENTER)
        setcell(r, 3, "{:,}".format(row["paid_amt"]), align=PP_ALIGN.RIGHT)
        setcell(r, 4, "{:,}".format(row["pipe"]), align=PP_ALIGN.RIGHT)
        for k in tot:
            tot[k] += row[k]
    lr = len(rows) + 1
    setcell(lr, 0, "TOTAL", True)
    setcell(lr, 1, str(tot["n"]), True, align=PP_ALIGN.CENTER)
    setcell(lr, 2, str(tot["paid"]), True, align=PP_ALIGN.CENTER)
    setcell(lr, 3, "{:,}".format(tot["paid_amt"]), True, align=PP_ALIGN.RIGHT)
    setcell(lr, 4, "{:,}".format(tot["pipe"]), True, align=PP_ALIGN.RIGHT)


def build_pptx(clients: list[Client], media_dir: str, out_path: str) -> None:
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H

    order = {s: i for i, s in enumerate(STATUSES)}
    ordered = sorted(clients, key=lambda c: (order.get(c.status, 0),
                                             c.slide_index))
    counts = status_counts(clients)
    subtitles = {
        "Pending approval": "Awaiting approval",
        "Approved": "Approved — payment pending",
        "Paid": "Paid — membership active",
    }
    current = None
    for client in ordered:
        if client.status != current:
            current = client.status
            _divider_slide(prs, current.upper(),
                           "%d member(s) · %s" %
                           (counts.get(current, 0), subtitles.get(current, "")))
        _member_slide(prs, client, media_dir)

    _divider_slide(prs, "ANALYSIS", "Membership analytics & totals")
    _analysis_slides(prs, clients)

    prs.save(out_path)
